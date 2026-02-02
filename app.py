import streamlit as st
from pptx import Presentation
from docx import Document
import io
import re

# --- ページ設定 ---
st.set_page_config(page_title="とも専用：レポート生成くん", layout="wide")
st.title("🚀 レポート生成アプリ (PPTX / DOCX)")

# --- サイドバー ---
st.sidebar.header("出力設定")
report_type = st.sidebar.radio("形式を選択", ["PowerPoint", "Word"])

# --- テキストエリア ---
st.subheader("1. Geminiの回答をここに貼り付け")
default_text = """# Slide 1: サンプル
- ここに内容を貼り付けてね"""
md_content = st.text_area("Markdownエディタ", value=default_text, height=400)

def clean_text(text):
    return re.sub(r'\*+', '', text)

# --- パワポ生成 ---
def create_pptx(content):
    prs = Presentation()
    raw_slides = re.split(r'\n#\s+', '\n' + content)
    for s in raw_slides:
        if not s.strip(): continue
        lines = s.strip().split('\n')
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = clean_text(lines[0])
        tf = slide.shapes.add_textbox(0, 0, 0, 0).text_frame # 簡易版
        tf.text = clean_text("\n".join(lines[1:]))
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

# --- Word生成 ---
def create_docx(content):
    doc = Document()
    raw_slides = re.split(r'\n#\s+', '\n' + content)
    for s in raw_slides:
        if not s.strip(): continue
        lines = s.strip().split('\n')
        doc.add_heading(clean_text(lines[0]), level=1)
        for line in lines[1:]:
            doc.add_paragraph(clean_text(line))
    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()

# --- 生成ボタン ---
if st.button("ファイルを生成"):
    if report_type == "PowerPoint":
        data = create_pptx(md_content)
        name = "report.pptx"
    else:
        data = create_docx(md_content)
        name = "report.docx"
    st.download_button("ダウンロード", data, file_name=name)
